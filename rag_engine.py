import os
from typing import List, Dict, Any, Optional
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import chromadb
# import camelot  # Moved to local import in _extract_tables
from chromadb.config import Settings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings
from langchain_pinecone import PineconeVectorStore
from pinecone import Pinecone

class RAGEngine:
    def __init__(self, persist_directory: str = "./chroma_db"):
        # Setup Cloud Vector DB (Pinecone) if keys are present
        self.pinecone_api_key = os.getenv("PINECONE_API_KEY")
        self.pinecone_index_name = os.getenv("PINECONE_INDEX_NAME")
        self.use_pinecone = bool(self.pinecone_api_key and self.pinecone_index_name)
        self._embeddings = None # Lazy load embeddings

        if self.use_pinecone:
            self.pc = Pinecone(api_key=self.pinecone_api_key)
            print("RAG Engine: Using Cloud Storage (Pinecone)")
        else:
            self.client = chromadb.PersistentClient(path=persist_directory)
            self.collection = self.client.get_or_create_collection(
                name="session_docs",
                metadata={"hnsw:space": "cosine"}
            )
            print("RAG Engine: Using Local Storage (ChromaDB)")

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )

    @property
    def embeddings(self):
        if self._embeddings is None:
            print("RAG Engine: Loading Embedding Model (FastEmbed)...")
            self._embeddings = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")
        return self._embeddings

    def _extract_text(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        if ext == ".pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in [".xlsx", ".xls"]:
            df = pd.read_excel(file_path)
            text = df.to_string()
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".json":
            with open(file_path, "r") as f:
                import json
                data = json.load(f)
                text = json.dumps(data, indent=2)
        elif ext in [".txt", ".md"]:
            with open(file_path, "r") as f:
                text = f.read()
        return text

    def _extract_tables(self, file_path: str) -> List[str]:
        import camelot # Local import to save memory if not used
        ext = os.path.splitext(file_path)[1].lower()
        table_markdowns = []
        if ext == ".pdf":
            try:
                # Limit pages to avoid memory exhaustion on large PDFs
                tables = camelot.read_pdf(file_path, pages='1-10', flavor='lattice')
                for table in tables:
                    table_markdowns.append(table.df.to_markdown(index=False))
            except Exception as e:
                try:
                    tables = camelot.read_pdf(file_path, pages='1-10', flavor='stream')
                    for table in tables:
                        table_markdowns.append(table.df.to_markdown(index=False))
                except Exception as e2:
                    print(f"Table extraction failed: {e2}")
        return table_markdowns

    def process_file(self, file_path: str, session_id: str, user_id: int, is_past_paper: bool = False):
        filename = os.path.basename(file_path)
        file_type = "past_paper" if is_past_paper else "document"
        
        # 1. Extract Text & Tables
        content = self._extract_text(file_path)
        tables = self._extract_tables(file_path)
        
        all_chunks = []
        all_metadatas = []
        
        if content.strip():
            chunks = self.text_splitter.split_text(content)
            all_chunks.extend(chunks)
            all_metadatas.extend([{"session_id": session_id, "user_id": user_id, "source": filename, "type": file_type} for _ in chunks])
            
        if tables:
            all_chunks.extend(tables)
            all_metadatas.extend([{"session_id": session_id, "user_id": user_id, "source": filename, "type": "table", "doc_type": file_type} for _ in tables])

        if not all_chunks:
            return

        # 2. Index in selected Vector DB
        if self.use_pinecone:
            # Pinecone VectorStore from_texts already handles batching internally or is less memory intensive on client
            PineconeVectorStore.from_texts(
                texts=all_chunks,
                embedding=self.embeddings,
                metadatas=all_metadatas,
                index_name=self.pinecone_index_name,
                namespace=f"user_{user_id}"
            )
        else:
            # Batching embeddings to avoid memory spikes
            batch_size = 10 # Reduced batch size for 512MB RAM
            for i in range(0, len(all_chunks), batch_size):
                batch_chunks = all_chunks[i : i + batch_size]
                batch_metadatas = all_metadatas[i : i + batch_size]
                
                vector_embeddings = self.embeddings.embed_documents(batch_chunks)
                ids = [f"u{user_id}_{session_id}_{filename}_{i+j}_{os.urandom(2).hex()}" for j in range(len(batch_chunks))]
                
                self.collection.add(
                    ids=ids,
                    embeddings=vector_embeddings,
                    documents=batch_chunks,
                    metadatas=batch_metadatas
                )

    def archive_chat_messages(self, session_id: str, user_id: int, messages: List[Dict[str, str]]):
        chunks = [f"{msg['role'].upper()}: {msg['content']}" for msg in messages]
        metadatas = [{"session_id": session_id, "user_id": user_id, "source": "Chat Archive", "type": "chat_archive"} for _ in chunks]
        
        if self.use_pinecone:
            PineconeVectorStore.from_texts(
                texts=chunks,
                embedding=self.embeddings,
                metadatas=metadatas,
                index_name=self.pinecone_index_name,
                namespace=f"user_{user_id}"
            )
        else:
            batch_size = 10
            for i in range(0, len(chunks), batch_size):
                batch_chunks = chunks[i : i + batch_size]
                batch_metadatas = metadatas[i : i + batch_size]
                
                ids = [f"archive_u{user_id}_{session_id}_{i+j}_{os.urandom(4).hex()}" for j in range(len(batch_chunks))]
                vector_embeddings = self.embeddings.embed_documents(batch_chunks)
                
                self.collection.add(
                    ids=ids, 
                    embeddings=vector_embeddings, 
                    documents=batch_chunks, 
                    metadatas=batch_metadatas
                )

    def query(self, session_id: str, question: str, user_id: int, top_k: int = 4, doc_type: Optional[str] = None) -> List[Dict[str, Any]]:
        formatted_results = []
        print(f"DEBUG: RAG Query for user {user_id}, session {session_id} - Question: {question}")
        
        if self.use_pinecone:
            vectorstore = PineconeVectorStore(index_name=self.pinecone_index_name, embedding=self.embeddings, namespace=f"user_{user_id}")
            # Filter by session_id and optionally doc_type
            filter_dict = {"session_id": session_id}
            if doc_type:
                filter_dict["type"] = doc_type
                
            results = vectorstore.similarity_search(question, k=top_k, filter=filter_dict)
            for doc in results:
                formatted_results.append({
                    "content": doc.page_content,
                    "source": doc.metadata.get("source", "Unknown")
                })
        else:
            query_embedding = self.embeddings.embed_query(question)
            
            # ChromaDB requires explicit $and for multiple conditions
            conditions = [
                {"session_id": session_id},
                {"user_id": user_id}
            ]
            if doc_type:
                conditions.append({"type": doc_type})
                
            where_clause = {"$and": conditions}

            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=top_k,
                where=where_clause
            )
            
            if results["documents"] and len(results["documents"][0]) > 0:
                print(f"DEBUG: RAG found {len(results['documents'][0])} results")
                for i in range(len(results["documents"][0])):
                    print(f"DEBUG: Result {i} Source: {results['metadatas'][0][i]['source']}")
                    formatted_results.append({
                        "content": results["documents"][0][i],
                        "source": results["metadatas"][0][i]["source"]
                    })
            else:
                print("DEBUG: RAG found NO results")
        return formatted_results

    def delete_session(self, session_id: str, user_id: int):
        """Wipes all documents and chat archives for a session/user."""
        if self.use_pinecone:
            vectorstore = PineconeVectorStore(index_name=self.pinecone_index_name, embedding=self.embeddings, namespace=f"user_{user_id}")
            # Pinecone deletion is more complex by metadata filter, normally done via namespace
            vectorstore.delete(delete_all=True)
        else:
            try:
                self.collection.delete(where={"$and": [{"session_id": session_id}, {"user_id": user_id}]})
            except Exception as e:
                print(f"Chroma delete failed (likely empty): {e}")

rag_engine = RAGEngine()
